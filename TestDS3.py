# not a malicious thing, just a bot to make fun of my friends!!
import os
os.environ["TOKENIZERS_PARALLELISM"] = "false"
import sys
# sys.modules["onnxruntime"] = type("onnxruntime", (), {"__spec__": type("spec", (), {"name": "onnxruntime"})()})()
import asyncio  # Добавить в начале файла
import pytesseract
from pdf2image import convert_from_path
from mcpsearch import MCPSearchFunction

from PIL import Image
import pandas as pd
from pptx import Presentation
from datetime import datetime
from pathlib import Path
import json  # <-- Add this line
import httpx
import chromadb
from chromadb.utils import embedding_functions
chromadb.utils.embedding_functions.DefaultEmbeddingFunction = lambda: None
from chromadb.config import Settings
from uuid import uuid4
import re
from chat_memory_manager import ChatMemoryManager
# from better_embedder import BetterEmbeddingFunction
from telegram import Update
import time
from telegram.ext import Application, CommandHandler, MessageHandler, filters, ContextTypes
from telegram import Update, Message, Chat, User
from pdfminer.high_level import extract_text  # Для PDF
from docx import Document  # Для DOCX
import requests

from PERSONAS import Persona, PERSONAS
from telegram import Update
import random
from telegram.ext import filters
from dotenv import load_dotenv
import logging
from collections import defaultdict, deque
from enum import Enum, auto
import base64
from io import BytesIO
from openai import OpenAI
from textblob import TextBlob
from persona_hooks import PERSONA_HOOKS
import sys
import warnings
warnings.filterwarnings("ignore", message="Failed to send telemetry event")

# chat_memories = defaultdict(lambda: deque(maxlen=32))
persona_contexts = defaultdict(dict)

load_dotenv()


def switch_persona(chat_id: int, user_id: int, new_persona: Persona) -> dict:
    key = (chat_id, user_id, new_persona.value)
    if key not in persona_contexts:
        persona_contexts[key] = {}
        persona_contexts[key]["memory_mgr"] = ChatMemoryManager(
            chat_id=chat_id,
            user_id=user_id,
            chroma_collection=chat_memory_collection,
            embedder=openai_embedder,
            openai_client=openai_client,
        )


    return persona_contexts[key]


SEARCH_TRIGGERS = ["search", "fetch", "lookup", "latest", "find", "news", "update", "найди", "поищи"]

searcher = MCPSearchFunction(
    google_api_key=os.getenv("GOOGLE_API_KEY"),
    search_engine_id=os.getenv("GOOGLE_CX_ID")
)

chroma_client = chromadb.PersistentClient(path="/data/chroma")



print("Chroma contents:", os.listdir("/data/chroma"))

print("📦 Checking existing Chroma collections...")
for col in chroma_client.list_collections():
    print(f"🧠 Found collection: {col.name}")

# Use OpenAI for embeddings
openai_embedder = embedding_functions.OpenAIEmbeddingFunction(
    api_key=os.getenv("OPENAI_API_KEY"),
    model_name="text-embedding-3-small"
)

logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

logging.getLogger("httpx").setLevel(logging.WARNING)
logging.getLogger("telegram").setLevel(logging.WARNING)

#local_embedder = BetterEmbeddingFunction("multilingual-e5-base")
# Initialize local embedder with error handling
# try:
#     local_embedder = BetterEmbeddingFunction("intfloat/multilingual-e5-base")
#     logger.info("✅ Local embedding model 'intfloat/multilingual-e5-base' loaded.")
#     _ = local_embedder.encode_documents(["trigger"])
#     logger.info("✅ Embedder preloaded with dummy input.")
# except Exception as e:
#     logger.critical(f"❌ Failed to load embedding model: {e}")
#     logger.critical("💡 Tip: Run 'python -c \"from sentence_transformers import SentenceTransformer; SentenceTransformer('intfloat/multilingual-e5-base')\"' to pre-download.")
#     exit(1)

# print(f"🔒 Embedder class: {local_embedder.__class__.__name__}")
# print(f"🧬 Model loaded: multilingual-e5-base")

def get_or_warn_collection(client, name, embedder):
    try:
        collection = client.get_collection(name=name)
        print(f"✅ Reusing existing Chroma collection: {name}")
        return collection
    except Exception as e:
        print(f"⚠️ Failed to get '{name}', creating new. Reason: {e}")
        return client.create_collection(name=name, embedding_function=embedder)

chroma_client.delete_collection("file_chunks")
chroma_client.delete_collection("chat_memory")


# Global collection: file chunks (we'll filter per chat later)
file_chunks_collection = chroma_client.get_or_create_collection(
    name="file_chunks",
    embedding_function=openai_embedder
)
# file_chunks_collection = get_or_warn_collection(chroma_client, "file_chunks", local_embedder)

# Collection for chat memory only (separate from files)
chat_memory_collection = chroma_client.get_or_create_collection(
    name="chat_memory",
    embedding_function=openai_embedder)
# chat_memory_collection = get_or_warn_collection(chroma_client, "chat_memory", local_embedder)


print("OPENAI_KEY_EXISTS:", "OPENAI_API_KEY" in os.environ)  # Debug line
openai_client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))

chat_modes = defaultdict(lambda: "normal")



# Initialize bot
app = Application.builder().token(os.getenv("BOT_TOKEN")).build()

# DeepSeek API setup
DEEPSEEK_API_URL = "https://api.deepseek.com/v1/chat/completions"
DEEPSEEK_HEADERS = {
    "Authorization": f"Bearer {os.getenv('DEEPSEEK_API_KEY')}",
    "Content-Type": "application/json"
}


# --------------------------------------
# EXACT FUNCTIONS FROM YOUR LIST (NO MORE, NO LESS)
# --------------------------------------


async def set_mode(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.message.chat.id
    user_id = update.effective_user.id  # <- ДОБАВЛЕНО для работы с persona_contexts

    if not context.args:
        await update.message.reply_text("Доступные режимы:\n" + "\n".join([f"- {p.value}" for p in Persona]))
        return

    mode_name = context.args[0].lower()
    try:
        persona = Persona(mode_name)
        # ДОБАВЛЕНО: Создаем/загружаем контекст новой персоны
        switch_persona(chat_id, user_id, persona)  # <- НОВАЯ СТРОКА
        chat_modes[chat_id] = persona.value
        await update.message.reply_text(f"🔹 Режим '{persona.value}' включён")
    except ValueError:
        await update.message.reply_text("❌ Неизвестный режим. Доступные: " + ", ".join([p.value for p in Persona]))


# --------------------------------------
# UTILITY FUNCTION
# --------------------------------------
def build_prompt(
        chat_id: int,
        user_input: str,
        persona_name: str,
        user_id: int = None,
        search_context: str = ""
) -> dict:
    persona = PERSONAS.get(Persona(persona_name), PERSONAS[Persona.NORMAL])
    context = persona_contexts.get(
        (chat_id, user_id, persona_name),
        {"message_history": deque(maxlen=32)}
    )

    # Format history with sender tags
    memory_mgr = context["memory_mgr"]

    parts = context["memory_mgr"].build_prompt_parts(user_input)
    history_str = "\n---\n".join(parts)

    return {
        "model": "gpt-4.1-mini",
        "messages": [
            {
                "role": "system",
                "content": (
                    f"{persona['system']}\n\n"
                    f"Если указаны АКТУАЛЬНЫЕ ДАННЫЕ, используй их в ответе.\n"
                    f"Добавляй строку 'Source' только если есть АКТУАЛЬНЫЕ ДАННЫЕ: ...' ТОЛЬКО с названием источника (например: РИА Новости, BBC, Интерфакс).\n"
                    f"Не добавляй ссылку (URL), кроме случаев, когда пользователь ЯВНО просит: 'дай ссылку', 'link', 'где источник' и т.п.\n\n"
                    f"АКТУАЛЬНЫЕ ДАННЫЕ:\n{search_context}\n\n"
                    f"ИСТОРИЯ:\n{history_str}"
                )
            },
            {
                "role": "user",
                "content": user_input
            }
        ],
        "temperature": persona["temperature"],
        "max_tokens": 900,
        "frequency_penalty": 1
    }


async def call_deepseek(payload: dict) -> str:
    """Call DeepSeek API with nuclear-grade quote prevention"""
    try:
        response = requests.post(
            DEEPSEEK_API_URL,
            headers=DEEPSEEK_HEADERS,
            json=payload,
            timeout=20  # Увеличенный таймаут
        )
        response.raise_for_status()

        raw_text = response.json()['choices'][0]['message']['content'].strip()

        # Strip edge quotes only
        if raw_text.startswith(('"', "'", "«")):
            raw_text = raw_text[1:]
        if raw_text.endswith(('"', "'", "»")):
            raw_text = raw_text[:-1]

        return raw_text or "Что-то не вышло. Давай еще."  # Сохраняем старый фолбек

    except requests.exceptions.RequestException as e:
        logger.error(f"API Error: {e}")
        if isinstance(e, requests.exceptions.Timeout):
            return "Сервер барахлит. Подожди минуту."  # Старая формулировка
        if response.status_code == 429:  # type: ignore
            return "Слишком много запросов."
        return "API не работает. Попробуй позже."

    except Exception as e:
        logger.critical(f"Critical: {e}")
        return "Поддержка уже работает над проблемой... Позвоните в OpenAI."  # Старый фолбек


async def call_ai(payload: dict) -> str:
    """DeepSeek with silent OpenAI fallback. Returns user-safe message if both fail."""
    TIMEOUT = 14

    try:
        # Try OpenAI first
        resp = openai_client.chat.completions.create(
            model="gpt-4.1-mini",  # Your specified model
            messages=payload["messages"],
            timeout=TIMEOUT,
        )
        return resp.choices[0].message.content

    except Exception as e:
        logger.error(f"OpenAI failed ({type(e).__name__}), falling back to DeepSeek")

        try:
            # Try DeepSeek
            response = requests.post(
                DEEPSEEK_API_URL,
                headers=DEEPSEEK_HEADERS,
                json=payload,
                timeout=TIMEOUT,
            )
            response.raise_for_status()  # Check HTTP errors (e.g., 4XX/5XX)
            return response.json()["choices"][0]["message"]["content"].strip()

        except Exception as fallback_error:
            logger.error(f"Both OpenAI and DeepSeek failed: {fallback_error}")
            return "Sorry, the AI service is temporarily unavailable. Please try again later."


# async def call_openai(input_text: str, system_prompt: str, temperature: float = 0.7, previous_response_id: str = None):
#     headers = {
#         "Authorization": f"Bearer {os.getenv('OPENAI_API_KEY')}",
#         "Content-Type": "application/json",
#     }

#     payload = {
#         "model": "gpt-4o-mini",
#         "input": [
#             {"role": "system", "content": system_prompt},
#             {"role": "user", "content": input_text},
#         ],
#         "temperature": temperature,
#         "store": True,
#         "tools": [{"type": "web_search_preview"}],
#     }

#     if previous_response_id:
#         payload["previous_response_id"] = previous_response_id

#     async with httpx.AsyncClient(timeout=20.0) as client:
#         response = await client.post(
#             "https://api.openai.com/v1/responses",
#             headers=headers,
#             json=payload,
#         )

#     response.raise_for_status()
#     data = response.json()

#     # Extract message text from the output
#     output = data.get("output", [])
#     message_obj = next((item for item in output if item.get("type") == "message"), None)
#     if not message_obj:
#         raise ValueError("No message found in response output")

#     content_list = message_obj.get("content", [])
#     text_entry = next((c for c in content_list if "text" in c), None)
#     if not text_entry:
#         raise ValueError("No text found in message content")

#     response_text = text_entry["text"]
#     response_id = data.get("id")

#     return response_text, response_id


async def handle_mention(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not update.message or not update.message.text:
        return

    chat_id = update.message.chat.id
    user_id = update.effective_user.id
    current_persona = Persona(chat_modes.get(chat_id, "normal"))
    text = update.message.text

    # 1. Load or create context
    persona_ctx = switch_persona(chat_id, user_id, current_persona)
    persona_ctx.setdefault("msg_counter", 0)
    persona_ctx["msg_counter"] += 1


    persona_ctx["memory_mgr"].add_message("user", text)

    text = update.message.text.strip().lower()

    search_context = ""
    if any(trigger in text for trigger in SEARCH_TRIGGERS):
        await update.message.reply_text("🌐 Searching...")
        result = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: searcher.enhanced_search(text)
        )
        search_context = result.get("combined_answer", "")

    logger.info("🔍 Injected Search Context:\n%s", search_context[:100])

    payload = build_prompt(
        chat_id=chat_id,
        user_input=text,
        persona_name=current_persona.value,
        user_id=user_id,
        search_context=search_context
    )

    response = await call_ai(payload)

    # 7. Store bot response
    # persona_ctx["message_history"].append({
    #     "text": response,
    #     "sender": "bot",
    #     "persona": current_persona.value
    # })

    persona_ctx["memory_mgr"].add_message("bot", response)

    await update.message.reply_text(response)


async def handle_reply(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not update.message.reply_to_message.from_user.id == context.bot.id:
        return

    chat_id = update.message.chat.id
    user_id = update.effective_user.id
    current_persona = Persona(chat_modes.get(chat_id, "normal"))

    # Use structured history instead of chat_memories
    persona_ctx = switch_persona(chat_id, user_id, current_persona)

    persona_ctx["memory_mgr"].add_message("user", update.message.text)

    payload = build_prompt(
        chat_id=chat_id,
        user_input=update.message.text,
        persona_name=current_persona.value,
        user_id=user_id
    )
    response = await call_ai(payload)
    await update.message.reply_text(response)

    persona_ctx["memory_mgr"].add_message("bot", response)


async def handle_file(update: Update, context: ContextTypes.DEFAULT_TYPE):
    # Initialize variables
    text = ""
    file_ext = os.path.splitext(update.message.document.file_name)[1].lower()

    # 1. File size check (20MB max)
    if update.message.document.file_size > 20_000_000:
        await update.message.reply_text("❌ Max 20MB")
        return

    # 2. Supported extensions
    ALLOWED_EXTENSIONS = [".pdf", ".docx", ".txt", ".csv", ".xlsx", ".pptx"]
    if file_ext not in ALLOWED_EXTENSIONS:
        await update.message.reply_text(f"❌ Only {', '.join(ALLOWED_EXTENSIONS)}")
        return

    # 3. Download file
    progress_msg = await update.message.reply_text("📥 Downloading...")
    file_path = f"/tmp/{int(time.time())}_{update.message.document.file_name}"
    try:
        telegram_file = await update.message.document.get_file()
        await telegram_file.download_to_drive(custom_path=file_path)

        # 4. Extract text (ALL FILE TYPES INCLUDED)
        if file_ext == ".pdf":
            text = extract_text(file_path)
            if not text.strip():
                images = convert_from_path(file_path)
                text = "\n".join(pytesseract.image_to_string(img, lang='rus+eng') for img in images)
        elif file_ext == ".docx":
            doc = Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text)
        elif file_ext == ".xlsx":
            df = pd.read_excel(file_path, engine='openpyxl')
            text = df.to_string()
        elif file_ext == ".pptx":
            prs = Presentation(file_path)
            text = "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))
        elif file_ext in (".txt", ".csv"):
            with open(file_path, "r", encoding="utf-8") as f:
                text = f.read()

        if not text.strip():
            await progress_msg.edit_text("🤷‍♂️ Empty/unreadable file")
            return

        # 5. Store chunks in Chroma
        chunks = split_text_into_chunks(text, chunk_size=800, overlap=100)
        store_chunks_in_chroma(
            chat_id=update.message.chat.id,
            user_id=update.effective_user.id,
            filename=update.message.document.file_name,
            chunks=chunks
        )

        # 6. Generate response with TIMEOUTS
        user_question = update.message.caption or "Резюме документа (5 предложений)"

        results = file_chunks_collection.query(
            query_texts=[user_question],
            n_results=3,
            where={"chat_id": str(update.message.chat.id)}
        )

        top_chunks = results.get("documents", [[]])[0]
        top_metadatas = results.get("metadatas", [[]])[0]

        # Include filenames as context prefix (only if they exist)
        file_labels = [meta.get("filename", "unknown") for meta in top_metadatas]
        file_info_text = "\n".join(f"[From file: {name}]" for name in set(file_labels))

        context_passage = file_info_text + "\n\n" + "\n\n".join(top_chunks)

        prompt = f"""
        Контекст из документа:
        {context_passage}

        Вопрос: {user_question}
        """

        response = None
        try:
            # Try DeepSeek with 15s timeout
            response = await asyncio.wait_for(
                call_deepseek({
                    "model": "deepseek-chat",
                    "messages": [{"role": "user", "content": prompt}],
                    "temperature": 0.3
                }),
                timeout=30.0
            )
        except (asyncio.TimeoutError, Exception):
            try:
                # Fallback to OpenAI with 10s timeout
                response = await asyncio.wait_for(
                    asyncio.get_event_loop().run_in_executor(
                        None,
                        lambda: openai_client.chat.completions.create(
                            model="gpt-4o-mini",
                            messages=[{"role": "user", "content": prompt}],
                            temperature=0.3,
                            timeout=20.0
                        ).choices[0].message.content
                    ),
                    timeout=12.0
                )
            except Exception as e:
                logger.error(f"AI Timeout: {e}")
                await progress_msg.edit_text("⌛ AI timeout. Try again.")
                return

        # 7. Update persona
        chat_id = update.message.chat.id
        persona_ctx = switch_persona(chat_id, update.effective_user.id, Persona(chat_modes[chat_id]))
        # persona_ctx["message_history"].append({
        #     "text": response,
        #     "sender": "bot",
        #     "persona": chat_modes[chat_id]
        # })

        persona_ctx["memory_mgr"].add_message("bot", response)
        
        await progress_msg.edit_text(response[:1200])

        context.user_data.update({
            'last_file_message_id': update.message.message_id,
            'is_file_context': True
        })

    except Exception as e:
        logger.error(f"FILE ERROR: {e}")
        await progress_msg.edit_text("💥 Processing error")

    finally:
        if os.path.exists(file_path):
            os.remove(file_path)


def split_text_into_chunks(text: str, chunk_size: int = 800, overlap: int = 100) -> list:
    """
    Splits text into overlapping chunks using simple sentence-based logic.

    Args:
        text (str): Full text to chunk.
        chunk_size (int): Approximate number of words per chunk.
        overlap (int): Words repeated between chunks for context.

    Returns:
        list[str]: List of text chunks.
    """
    sentences = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    current_chunk = []

    for sentence in sentences:
        words = sentence.split()
        current_chunk.extend(words)

        if len(current_chunk) >= chunk_size:
            chunks.append(' '.join(current_chunk[:chunk_size]))
            current_chunk = current_chunk[chunk_size - overlap:]

    if current_chunk:
        chunks.append(' '.join(current_chunk))

    return chunks


def store_chunks_in_chroma(chat_id, user_id, filename, chunks: list[str]):
    """
    Stores text chunks in Chroma with metadata.
    Args:
        chat_id (int): Telegram chat ID.
        user_id (int): Telegram user ID.
        filename (str): Original file name.
        chunks (list[str]): List of text chunks.
    """
    ids = [str(uuid4()) for _ in chunks]
    metadatas = [
        {
            "chat_id": str(chat_id),
            "user_id": str(user_id),
            "filename": filename,
            "chunk_index": i,
            "total_chunks": len(chunks)
        }
        for i in range(len(chunks))
    ]
    BATCH_SIZE = 100
    for i in range(0, len(chunks), BATCH_SIZE):
        file_chunks_collection.add(
            documents=chunks[i:i + BATCH_SIZE],
            metadatas=metadatas[i:i + BATCH_SIZE],
            ids=ids[i:i + BATCH_SIZE]
        )
    logger.info(f"✅ Stored {len(chunks)} chunks for '{filename}' (chunk indices 0–{len(chunks)-1})")


async def handle_file_query(update: Update, context: ContextTypes.DEFAULT_TYPE):
    # 1. Check triggers
    if not update.message or not update.message.text:
        return
    query = update.message.text.lower()
    is_reply_to_file_summary = (
            update.message.reply_to_message and
            update.message.reply_to_message.from_user.id == context.bot.id and
            update.message.reply_to_message.message_id == context.user_data.get('last_file_message_id')
    )

    if not any(trigger in query for trigger in
               ["файл", "в файле", "files", "file", "document"]) and not is_reply_to_file_summary:
        logger.debug(f"Not a file query: '{query}' (reply_to_file={is_reply_to_file_summary})")
        await handle_mention(update, context)
        return

    # 3. Query Chroma for top chunks
    query = update.message.text
    # Debug: Show all chunks for this chat_id

    results = file_chunks_collection.query(
        query_texts=[query],
        n_results=3,
        where={"chat_id": str(update.message.chat.id)}
    )

    logger.info("📄 FILE QUERY TRIGGERED via Chroma")
    logger.info(f"User query: {query}")
    logger.info(f"Trigger: {'keywords' if not is_reply_to_file_summary else 'reply to file summary'}")

    top_chunks = results.get("documents", [[]])[0]
    if not top_chunks:
        logger.info("⚠️ No chunks returned from Chroma query")
        await update.message.reply_text("⚠️ Не нашёл ничего в документе по запросу.")
        return
    logger.info(f"Top Chunks count: {len(top_chunks)}")
    context_passage = "\n\n".join(top_chunks)

    full_prompt = f"""
    Контекст из документа:
    {context_passage}

    Вопрос: {query}
    """

    # 4. Process with AI (minimal system prompt)
    try:
        prompt_payload = {
            "model": "deepseek-chat",
            "messages": [{
                "role": "system",
                "content": "Используй ТОЛЬКО текст из документа ниже. Игнорируй внешние знания."
            }, {
                "role": "user",
                "content": full_prompt
            }],
            "temperature": 0.3
        }

        try:
            # Try DeepSeek first (with 15s timeout)
            response = await asyncio.wait_for(
                call_deepseek(prompt_payload),
                timeout=30.0
            )
        except (asyncio.TimeoutError, Exception):
            logger.warning("DeepSeek failed, falling back to OpenAI")
            # Fallback to OpenAI with timeout
            response = await asyncio.wait_for(
                asyncio.get_event_loop().run_in_executor(
                    None,
                    lambda: openai_client.chat.completions.create(
                        model="gpt-4o-mini",
                        messages=prompt_payload["messages"],
                        temperature=0.3,
                        timeout=20.0
                    ).choices[0].message.content
                ),
                timeout=12.0
            )

        persona_ctx = switch_persona(
            chat_id=update.message.chat.id,
            user_id=update.effective_user.id,
            new_persona=Persona(chat_modes[update.message.chat.id])
        )
        
        persona_ctx["memory_mgr"].add_message("bot", response)
        
        await update.message.reply_text(response[:1200])

    except Exception as e:
        logger.error(f"File query failed after fallback: {e}")
        await update.message.reply_text("💥 Ошибка обработки запроса")


async def handle_image(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if not update.message.photo:
        await update.message.reply_text("Это не изображение.")
        return

    try:
        # --- Original Image Processing (UNTOUCHED) ---
        photo_file = await update.message.photo[-1].get_file()
        photo_bytes = BytesIO()
        await photo_file.download_to_memory(out=photo_bytes)
        base64_image = base64.b64encode(photo_bytes.getvalue()).decode('utf-8')

        user_question = (
                update.message.caption or
                (update.message.reply_to_message.text if update.message.reply_to_message else None) or
                "Опиши это изображение. "
        )

        # --- NEW: Structured Memory Integration ---
        chat_id = update.message.chat.id
        user_id = update.effective_user.id
        current_persona = Persona(chat_modes.get(chat_id, "normal"))  # Your original persona logic
        persona_ctx = switch_persona(chat_id, user_id, current_persona)

        # Store image event (formatted to match your style)
        # persona_ctx["message_history"].append({
        #     "text": f"[Image: {user_question}]",
        #     "sender": "user",
        #     "persona": None
        # })
        persona_ctx["memory_mgr"].add_message("user", f"[Image: {user_question}]")

        
        # --- Your Original API Call (EXACTLY AS IS) ---
        persona_config = PERSONAS[current_persona]
        prompt_text = (
            f"{persona_config['system']}\n\n"
            f"Запрос: {user_question}\n\n"
            "Ответь в своём стиле (макс. 10 предложений). Удели мнимание деталям изображения."
        )

        processing_msg = await update.message.reply_text("Разглядываю")
        response = openai_client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[{
                "role": "user",
                "content": [
                    {"type": "text", "text": prompt_text},
                    {
                        "type": "image_url",
                        "image_url": {
                            "url": f"data:image/jpeg;base64,{base64_image}",
                            "detail": "low"  # Your original detail level
                        }
                    }
                ]
            }],
            max_tokens=800
        )

        # --- NEW: Store Bot Response ---
        analysis = response.choices[0].message.content
        # persona_ctx["message_history"].append({
        #     "text": analysis,
        #     "sender": "bot",
        #     "persona": current_persona.value
        # })

        persona_ctx["memory_mgr"].add_message("bot", analysis)

        
        await processing_msg.edit_text(analysis[:1000])  # Your original truncation

    except Exception as e:
        logger.error(f"Image error: {e}")
        await update.message.reply_text("Не вышло. Попробуй другую картинку.")


async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    try:
        # 1. Your existing transcription code
        voice_file = await update.message.voice.get_file()
        voice_bytes = BytesIO()
        await voice_file.download_to_memory(out=voice_bytes)
        user_text = openai_client.audio.transcriptions.create(
            model="whisper-1",
            file=("voice.ogg", voice_bytes.getvalue())
        ).text.strip()

        if not user_text:
            await update.message.reply_text("🔇 Пустое сообщение")
            return

        # 2. ONLY NEW PART: Check for file queries
        if any(trigger in user_text.lower() for trigger in ["файл", "в файле", "документ", "в документе", "document"]):
            results = file_chunks_collection.query(
                query_texts=[user_text],
                n_results=3,
                where={"chat_id": str(update.message.chat.id)}
            )

            top_chunks = results.get("documents", [[]])[0]
            context_passage = "\n\n".join(top_chunks)

            full_prompt = f"""
            Контекст из документа:
            {context_passage}

            Вопрос: {user_text}
            """

            response = await call_ai({
                "model": "deepseek-chat",
                "messages": [{
                    "role": "system",
                    "content": "Отвечай ТОЛЬКО из файла. Не используй внешние знания."
                }, {
                    "role": "user",
                    "content": full_prompt
                }],
                "temperature": 0.3
            })
            await update.message.reply_text(response[:1200])
            return  # Important! Skip normal handling

        # 3. Original behavior for non-file queries
        msg = update.message
        msg._unfreeze()
        msg.text = user_text
        msg.voice = None
        msg._freeze()

        handler = handle_mention if msg.chat.type == "private" else group_handler
        await handler(update, context)

    except Exception as e:
        logger.error(f"Voice error: {e}")
        await update.message.reply_text("Error processing voice")

async def list_files(update: Update, context: ContextTypes.DEFAULT_TYPE):
    chat_id = update.message.chat.id
    user_id = str(update.effective_user.id)

    # Query Chroma for all unique filenames from this user+chat
    results = file_chunks_collection.get(
        where={"$and": [
            {"chat_id": {"$eq": str(chat_id)}},
            {"user_id": {"$eq": str(user_id)}}
        ]}
    )
    if not results["documents"]:
        await update.message.reply_text("No files found.")
        return

    # Extract unique filenames
    filenames = {meta["filename"] for meta in results["metadatas"]}
    file_list = "\n".join(f"📄 {name}" for name in sorted(filenames))
    await update.message.reply_text(f"Your files:\n{file_list}")



async def group_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    # Check for either @mention OR reply-to-bot
    bot_username = context.bot.username.lower()
    is_reply_to_bot = (
            update.message.reply_to_message and
            update.message.reply_to_message.from_user.id == context.bot.id
    )
    has_mention = (
            (update.message.text and f"@{bot_username}" in update.message.text.lower()) or
            (update.message.caption and f"@{bot_username}" in update.message.caption.lower())
    )

    if not (is_reply_to_bot or has_mention):
        return  # Ignore normal group messages

    chat_id = update.message.chat.id
    user_id = update.effective_user.id

    # Route to appropriate handler
    if update.message.photo:
        await handle_image(update, context)
    elif update.message.document:
        await handle_file(update, context)
    else:
        # === ЕДИНСТВЕННОЕ ИЗМЕНЕНИЕ ===
        payload = build_prompt(
            chat_id=chat_id,
            user_input=update.message.text,
            persona_name=chat_modes[chat_id],
            user_id=user_id)

        response = await call_ai(payload)
        await update.message.reply_text(response)
    # --------------------------------------


# REGISTER ALL COMMANDS
# --------------------------------------
commands = [

    ("mode", set_mode),
    ('files', list_files)

]

for cmd, handler in commands:
    app.add_handler(CommandHandler(cmd, handler))

# ===== 2. GROUP CHAT HANDLER (REPLACE with this) =====
app.add_handler(MessageHandler(
    filters.ChatType.GROUPS & (filters.TEXT | filters.PHOTO | filters.Document.ALL),
    group_handler  # Your custom function that checks @mentions
))

# file query handler
app.add_handler(MessageHandler(
    filters.TEXT & ~filters.COMMAND,
    handle_file_query  # <-- this function will internally decide
))

# ===== 1. PRIVATE CHAT HANDLER (keep) =====
app.add_handler(MessageHandler(
    filters.ChatType.PRIVATE & (filters.TEXT | filters.PHOTO | filters.Document.ALL),
    lambda update, ctx: (
        handle_image(update, ctx) if update.message.photo else
        handle_file(update, ctx) if update.message.document else
        handle_mention(update, ctx)
    )
))

# ===== 3. REPLY HANDLER (keep one) =====
app.add_handler(MessageHandler(
    filters.TEXT & filters.REPLY,
    handle_reply
))

app.add_handler(MessageHandler(filters.VOICE, handle_voice))
app.add_handler(CommandHandler("files", list_files))




if __name__ == "__main__":
    print("New TestHelper launched")
    app.run_polling()







